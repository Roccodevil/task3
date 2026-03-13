import fitz  # PyMuPDF
from pptx import Presentation
import pandas as pd
import os
import config
from typing import Tuple, List
from pathlib import Path


def process_document(file_path: str) -> Tuple[str, List[str]]:
    """
    Identifies file type, extracts text for Ollama, and saves images locally.
    """
    file_ext = Path(file_path).suffix.lower().lstrip(".")
    extracted_text = ""
    saved_image_paths = []

    print(f"--- PARSING FILE: {file_path} ---")

    # Ensure temp dir exists
    os.makedirs(config.TEMP_IMG_DIR, exist_ok=True)

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")

    # 1. Handle PDF
    if file_ext == "pdf":
        try:
            with fitz.open(file_path) as doc:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    extracted_text += page.get_text() + "\n"

                    # Extract PDF Images
                    for img_index, img in enumerate(page.get_images(full=True)):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image.get("image")
                        image_ext = base_image.get("ext", "png")
                        if not image_bytes:
                            continue
                        img_filename = os.path.join(config.TEMP_IMG_DIR, f"pdf_p{page_num}_i{img_index}.{image_ext}")

                        with open(img_filename, "wb") as f:
                            f.write(image_bytes)
                        saved_image_paths.append(img_filename)
        except Exception as e:
            raise ValueError(f"Error parsing PDF: {e}") from e

    # 2. Handle PowerPoint (PPTX)
    elif file_ext == "pptx":
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    # Extract Text
                    if hasattr(shape, "text") and shape.text:
                        extracted_text += shape.text + "\n"

                    # Extract Images
                    try:
                        if hasattr(shape, "image") and shape.image is not None:
                            image = shape.image
                            image_name = image.filename or f"image_{slide_num}.png"
                            image_name = os.path.basename(image_name)
                            img_filename = os.path.join(config.TEMP_IMG_DIR, f"pptx_s{slide_num}_{image_name}")
                            with open(img_filename, "wb") as f:
                                f.write(image.blob)
                            saved_image_paths.append(img_filename)
                    except Exception:
                        # Some shapes may not expose image attributes uniformly
                        continue
        except Exception as e:
            raise ValueError(f"Error parsing PPTX: {e}") from e

    # 3. Handle Spreadsheets / CSV
    elif file_ext in ["csv", "xlsx", "xls"]:
        try:
            if file_ext == "csv":
                try:
                    df = pd.read_csv(file_path)
                except UnicodeDecodeError:
                    df = pd.read_csv(file_path, encoding="latin1")
            else:
                df = pd.read_excel(file_path)

            extracted_text += "Extracted Data Table:\n"
            try:
                extracted_text += df.to_markdown(index=False) + "\n"
            except Exception:
                extracted_text += df.to_string(index=False) + "\n"
        except Exception as e:
            raise ValueError(f"Error parsing table file: {e}") from e

    # 4. Handle Raw Images
    elif file_ext in ["jpg", "jpeg", "png","webp","avif"]:
        extracted_text += "Image file provided. Handing off to Vision model."
        saved_image_paths.append(file_path)

    else:
        raise ValueError(f"Unsupported format: {file_ext}")

    if not extracted_text:
        extracted_text = "No textual content extracted from input."

    return extracted_text, saved_image_paths
